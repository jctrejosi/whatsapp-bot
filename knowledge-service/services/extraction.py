"""Multi-format text extraction service.

Supports: PDF, DOCX, PPTX, XLSX, XLS, ODT, ODS, ODP, RTF,
          CSV, TXT, MD, HTML, JSON, XML, YAML, and code files (.py, .js, .ts, .java, .cpp, etc.)
"""

import csv
import io
import os
import re
import zipfile
import xml.etree.ElementTree as ET
from typing import NamedTuple


class ExtractionResult(NamedTuple):
    """Result of extracting text from a file."""
    pages: list[dict]          # [{page, text, char_count}, ...]
    full_text: str
    total_pages: int
    total_chars: int
    metadata: dict | None


# ─── File type helpers ───────────────────────────────────────────────────────────

# Extensions that are plain-text and can be read directly
PLAIN_TEXT_EXTENSIONS = {
    '.txt', '.md', '.rst', '.log',
    '.py', '.js', '.ts', '.jsx', '.tsx', '.java', '.cpp', '.c', '.h', '.hpp',
    '.rb', '.go', '.rs', '.swift', '.kt', '.scala', '.php', '.lua', '.r',
    '.sh', '.bash', '.zsh', '.fish', '.ps1', '.bat',
    '.html', '.htm', '.css', '.scss', '.less',
    '.svg', '.xml', '.json', '.yaml', '.yml', '.toml', '.ini', '.cfg',
    '.url', '.tex', '.sql', '.graphql',
}

# Structured plain-text that should be read field-by-field
STRUCTURED_TEXT_EXTENSIONS = {'.csv', '.tsv'}

# Modern Office formats (ZIP + XML)
DOCX_EXTENSIONS = {'.docx'}
PPTX_EXTENSIONS = {'.pptx'}
XLSX_EXTENSIONS = {'.xlsx'}
ODF_EXTENSIONS = {'.odt', '.ods', '.odp'}  # OpenDocument (LibreOffice)

# Legacy binary Office formats — try plain text (rarely works, but won't crash)
LEGACY_BINARY_EXTENSIONS = {'.doc', '.ppt', '.xls'}

# RTF
RTF_EXTENSIONS = {'.rtf'}

# Compressed archives
ARCHIVE_EXTENSIONS = {'.zip'}


def extract_file_bytes(data: bytes, filename: str) -> ExtractionResult:
    """Extract text from uploaded file bytes, auto-detecting the format."""
    ext = os.path.splitext(filename)[1].lower()

    if ext == '.pdf':
        return _extract_pdf_bytes(data)

    if ext in DOCX_EXTENSIONS:
        return _extract_docx(data)

    if ext in PPTX_EXTENSIONS:
        return _extract_pptx(data)

    if ext in XLSX_EXTENSIONS:
        return _extract_xlsx(data)

    if ext == '.xls':
        return _extract_xls(data)

    if ext in ODF_EXTENSIONS:
        return _extract_odf(data)

    if ext in RTF_EXTENSIONS:
        return _extract_rtf(data)

    if ext in ARCHIVE_EXTENSIONS:
        return _extract_archive(data, filename)

    if ext in STRUCTURED_TEXT_EXTENSIONS:
        return _extract_csv_text(data, filename)

    if ext in LEGACY_BINARY_EXTENSIONS:
        return _extract_plain_text(data, filename, note='formato binario legado — texto crudo')

    # All other text-based formats (code, data, web, etc.)
    return _extract_plain_text(data, filename)


# ─── PDF: pipeline inteligente que elige el mejor extractor ─────────────

def _extract_pdf_bytes(data: bytes) -> ExtractionResult:
    """Extract text from PDF using the best available method.
    Tries PyMuPDF first (fast, good for most). Falls back to pdfplumber
    for table-heavy PDFs, and Tesseract OCR for scanned/image PDFs."""

    # 1. PyMuPDF (word-level) — primary
    result_pymupdf = _extract_pdf_pymupdf(data)
    score_pymupdf = _quality_score(result_pymupdf)

    # 2. pdfplumber — fallback for table-heavy or complex layouts
    result_plumber = None
    score_plumber = 0
    if score_pymupdf < 0.6:  # calidad baja, intentar pdfplumber
        try:
            result_plumber = _extract_pdf_plumber(data)
            score_plumber = _quality_score(result_plumber)
        except Exception:
            pass

    # 3. OCR (Tesseract) — fallback for scanned/image PDFs
    result_ocr = None
    score_ocr = 0
    if score_pymupdf < 0.2 and score_plumber < 0.2:  # probablemente escaneado
        try:
            result_ocr = _extract_pdf_ocr(data)
            score_ocr = _quality_score(result_ocr)
        except Exception:
            pass

    # Elegir el mejor
    best = result_pymupdf
    best_score = score_pymupdf
    best_name = 'pymupdf'
    if result_plumber and score_plumber > best_score:
        best = result_plumber
        best_score = score_plumber
        best_name = 'pdfplumber'
    if result_ocr and score_ocr > best_score:
        best = result_ocr
        best_score = score_ocr
        best_name = 'ocr'

    print(f'  📄 PDF extraction: pymupdf={score_pymupdf:.2f} plumber={score_plumber:.2f} ocr={score_ocr:.2f} → using {best_name}')
    return best


def _quality_score(result: ExtractionResult) -> float:
    """Rate extraction quality 0-1 based on chars/page and text cleanliness."""
    pages = result.total_pages or 1
    chars_per_page = result.total_chars / pages

    # Ideal: ~1500-3000 chars/page. Below 50 suggests scanned/image PDF.
    if chars_per_page > 2000:
        density = 1.0
    elif chars_per_page > 500:
        density = 0.6 + 0.4 * (chars_per_page - 500) / 1500
    elif chars_per_page > 50:
        density = 0.3 * (chars_per_page - 50) / 450
    else:
        density = 0.0

    # Penalizar texto con muchos duplicados ("DELUXEDELUXE")
    text = result.full_text
    if len(text) > 100:
        words = text.split()
        if len(words) > 20:
            dup_ratio = sum(1 for i in range(1, len(words)) if words[i] == words[i-1]) / len(words)
            # 0% duplicates = no penalty, 30%+ = heavy penalty
            dup_penalty = min(1.0, dup_ratio / 0.3)
            density *= (1.0 - 0.5 * dup_penalty)

    return max(0.0, min(1.0, density))


# ─── PyMuPDF (word-level) — mejor para texto y layouts complejos ─────────

def _extract_pdf_pymupdf(data: bytes) -> ExtractionResult:
    """Extract text from PDF using PyMuPDF word-level extraction with spatial sorting."""
    import fitz
    doc = fitz.open(stream=data, filetype='pdf')
    pages_out = []
    parts = []

    try:
        for i, page in enumerate(doc):
            words = page.get_text('words')

            if not words:
                continue

            # Agrupar palabras en líneas por coordenada y
            tolerance = 4
            lines = []
            current_line = []
            current_y = None

            for w in sorted(words, key=lambda w: (round(w[1]), w[0])):
                wy = w[1]
                if current_y is None:
                    current_y = wy
                    current_line.append(w)
                elif abs(wy - current_y) < tolerance:
                    current_line.append(w)
                else:
                    if current_line:
                        lines.append(current_line)
                    current_line = [w]
                    current_y = wy
            if current_line:
                lines.append(current_line)

            # Reconstruir texto línea por línea
            page_lines = []
            for line_words in lines:
                sorted_words = sorted(line_words, key=lambda w: w[0])
                text_parts = []
                prev = ''
                for w in sorted_words:
                    word = w[4]
                    if word != prev:
                        text_parts.append(word)
                    prev = word
                line_text = ' '.join(text_parts).strip()
                if line_text:
                    page_lines.append(line_text)

            # Detectar párrafos por gaps verticales
            paragraphs = []
            current_para = []
            prev_y = None
            para_gap = 10

            for line_words, line_text in zip(lines, page_lines):
                if not line_text:
                    continue
                y = line_words[0][1]
                if prev_y is not None and (y - prev_y) > para_gap:
                    if current_para:
                        paragraphs.append('\n'.join(current_para))
                    current_para = []
                current_para.append(line_text)
                prev_y = y
            if current_para:
                paragraphs.append('\n'.join(current_para))

            page_text = '\n\n'.join(paragraphs)

            if page_text:
                pages_out.append({'page': i + 1, 'text': page_text, 'char_count': len(page_text)})
                parts.append(page_text)

        meta = doc.metadata
        metadata = {
            'title': meta.get('title'),
            'author': meta.get('author'),
            'creator': meta.get('creator'),
            'producer': meta.get('producer'),
            'format': meta.get('format'),
        } if meta else None

    finally:
        doc.close()

    full_text = '\n\n'.join(parts)
    return ExtractionResult(
        pages=pages_out,
        full_text=full_text,
        total_pages=len(pages_out),
        total_chars=len(full_text),
        metadata=metadata,
    )


# ─── pdfplumber — mejor para tablas en PDFs con markup ───────────────────

def _extract_pdf_plumber(data: bytes) -> ExtractionResult:
    """Extract text using pdfplumber — excels at table-heavy PDFs."""
    import pdfplumber, io
    pages_out = []
    parts = []

    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for i, page in enumerate(pdf.pages):
            # Intentar extraer tablas primero
            tables = page.extract_tables()
            table_texts = []
            for table in tables:
                if not table:
                    continue
                # Formatear tabla como texto alineado
                rows = []
                for row in table:
                    if row:
                        rows.append(' | '.join(str(c or '') for c in row))
                if rows:
                    table_texts.append('\n'.join(rows))

            # Extraer texto normal
            text = (page.extract_text() or '').strip()

            # Combinar: tablas primero, luego texto
            combined = []
            if table_texts:
                combined.append('\n\n'.join(table_texts))
            if text:
                # Evitar duplicar contenido de tablas en el texto
                # (pdfplumber incluye texto de tablas en extract_text también)
                combined.append(text)

            page_text = '\n\n'.join(combined).strip()
            if page_text:
                pages_out.append({'page': i + 1, 'text': page_text, 'char_count': len(page_text)})
                parts.append(page_text)

    full_text = '\n\n'.join(parts)
    return ExtractionResult(
        pages=pages_out,
        full_text=full_text,
        total_pages=len(pages_out),
        total_chars=len(full_text),
        metadata=None,
    )


# ─── OCR (Tesseract) — para PDFs escaneados/imagen ───────────────────────

def _extract_pdf_ocr(data: bytes) -> ExtractionResult:
    """Extract text from scanned/image PDFs using Tesseract OCR."""
    import fitz, io
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        raise RuntimeError('pytesseract or Pillow not installed — cannot OCR PDF')

    doc = fitz.open(stream=data, filetype='pdf')
    pages_out = []
    parts = []

    try:
        for i, page in enumerate(doc):
            # Renderizar página como imagen a 200 DPI
            pix = page.get_pixmap(dpi=200)
            img = Image.open(io.BytesIO(pix.tobytes('png')))
            text = pytesseract.image_to_string(img, lang='eng+spa').strip()
            if text:
                pages_out.append({'page': i + 1, 'text': text, 'char_count': len(text)})
                parts.append(text)
    finally:
        doc.close()

    full_text = '\n\n'.join(parts)
    return ExtractionResult(
        pages=pages_out,
        full_text=full_text,
        total_pages=len(pages_out),
        total_chars=len(full_text),
        metadata=None,
    )


# ─── DOCX (python-docx) ────────────────────────────────────────────────────────

def _extract_docx(data: bytes) -> ExtractionResult:
    from docx import Document
    doc = Document(io.BytesIO(data))
    parts = []
    for p in doc.paragraphs:
        t = p.text.strip()
        if t:
            parts.append(t)
    # Extract table text as well
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(' | '.join(cells))
    full_text = '\n\n'.join(parts)
    return ExtractionResult(
        pages=[{'page': 1, 'text': full_text, 'char_count': len(full_text)}],
        full_text=full_text,
        total_pages=1,
        total_chars=len(full_text),
        metadata=None,
    )


# ─── PPTX (python-pptx) ────────────────────────────────────────────────────────

def _extract_pptx(data: bytes) -> ExtractionResult:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    t = paragraph.text.strip()
                    if t:
                        texts.append(t)
        txt = '\n'.join(texts)
        if txt.strip():
            slides.append({'page': i + 1, 'text': txt, 'char_count': len(txt)})
    full_text = '\n\n---\n\n'.join(s.get('text', '') for s in slides)
    return ExtractionResult(
        pages=slides,
        full_text=full_text,
        total_pages=len(slides),
        total_chars=len(full_text),
        metadata=None,
    )


# ─── XLSX (openpyxl) ───────────────────────────────────────────────────────────

def _extract_xlsx(data: bytes) -> ExtractionResult:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    sheets = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows_text = []
        for row in ws.iter_rows(values_only=True):
            row_text = ' | '.join(str(c) if c is not None else '' for c in row).strip(' |')
            if row_text:
                rows_text.append(row_text)
        if rows_text:
            sheets.append(f'--- Hoja: {sheet_name} ---\n' + '\n'.join(rows_text))
    full_text = '\n\n'.join(sheets)
    return ExtractionResult(
        pages=[{'page': 1, 'text': full_text, 'char_count': len(full_text)}],
        full_text=full_text,
        total_pages=1,
        total_chars=len(full_text),
        metadata=None,
    )


# ─── XLS (xlrd) — legacy binary ────────────────────────────────────────────────

def _extract_xls(data: bytes) -> ExtractionResult:
    try:
        import xlrd
        wb = xlrd.open_workbook(file_contents=data)
        sheets = []
        for sheet in wb.sheets():
            rows_text = []
            for row_idx in range(sheet.nrows):
                row_text = ' | '.join(
                    str(sheet.cell_value(row_idx, col_idx)) for col_idx in range(sheet.ncols)
                ).strip(' |')
                if row_text:
                    rows_text.append(row_text)
            if rows_text:
                sheets.append(f'--- Hoja: {sheet.name} ---\n' + '\n'.join(rows_text))
        full_text = '\n\n'.join(sheets) or _extract_plain_text(data, '.xls', 'sin contenido estructurado').full_text
    except Exception:
        return _extract_plain_text(data, '.xls', 'no se pudo parsear como XLS binario')
    return ExtractionResult(
        pages=[{'page': 1, 'text': full_text, 'char_count': len(full_text)}],
        full_text=full_text,
        total_pages=1,
        total_chars=len(full_text),
        metadata=None,
    )


# ─── ODF (LibreOffice: ODT, ODS, ODP) — ZIP + XML ──────────────────────────────

def _extract_odf(data: bytes) -> ExtractionResult:
    """Extract text from OpenDocument files using XML namespaces."""
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            if 'content.xml' not in zf.namelist():
                return _extract_plain_text(data, '.odf', 'no es un ODF válido')
            root = ET.fromstring(zf.read('content.xml'))

        # OpenDocument namespaces
        ns = {
            'text': 'urn:oasis:names:tc:opendocument:xmlns:text:1.0',
            'table': 'urn:oasis:names:tc:opendocument:xmlns:table:1.0',
            'office': 'urn:oasis:names:tc:opendocument:xmlns:office:1.0',
        }
        # Try multiple namespace variants (some ODF files use different prefixes)
        # Just strip all tags and keep text content from the body
        paragraphs = []
        for p in root.iter():
            tag = p.tag.split('}')[-1] if '}' in p.tag else p.tag
            if tag in ('p', 'h'):
                txt = ''.join(p.itertext()).strip()
                if txt:
                    paragraphs.append(txt)
            elif tag == 'tab':
                paragraphs.append('\t')
        text = '\n\n'.join(paragraphs)
        if not text:
            # Fallback: strip all tags
            content = ET.tostring(root, encoding='unicode')
            text = re.sub(r'<[^>]+>', ' ', content)
            text = re.sub(r'\s+', ' ', text).strip()
    except Exception:
        return _extract_plain_text(data, '.odf', 'no se pudo extraer de ODF')
    return ExtractionResult(
        pages=[{'page': 1, 'text': text, 'char_count': len(text)}],
        full_text=text,
        total_pages=1,
        total_chars=len(text),
        metadata=None,
    )


# ─── RTF ───────────────────────────────────────────────────────────────────────

def _extract_rtf(data: bytes) -> ExtractionResult:
    """Robust RTF text extractor — handles control words, hex escapes, unicode, and font/color tables."""
    text = data.decode('utf-8', errors='replace')

    # Remove RTF groups that don't contain readable text
    text = re.sub(r'\{\\*?\\fonttbl.*?\}', ' ', text, flags=re.DOTALL)
    text = re.sub(r'\{\\*?\\colortbl.*?\}', ' ', text, flags=re.DOTALL)
    text = re.sub(r'\{\\*?\\stylesheet.*?\}', ' ', text, flags=re.DOTALL)
    text = re.sub(r'\{\\*?\\info.*?\}', ' ', text, flags=re.DOTALL)
    text = re.sub(r'\{\\*?\\pntext.*?\}', ' ', text, flags=re.DOTALL)
    text = re.sub(r'\{\\*?\\listtext.*?\}', ' ', text, flags=re.DOTALL)
    text = re.sub(r'\{\\*?\\xmlnstbl.*?\}', ' ', text, flags=re.DOTALL)

    # Handle hex-escaped characters: \'e1 → á
    def _hex_replace(m):
        try:
            return bytes([int(m.group(1), 16)]).decode('cp1252', errors='replace')
        except Exception:
            return ' '
    text = re.sub(r"\\'([0-9a-fA-F]{2})", _hex_replace, text)

    # Handle unicode escapes: \u1234? → unicode char
    def _unicode_replace(m):
        try:
            return chr(int(m.group(1)))
        except Exception:
            return ' '
    text = re.sub(r'\\u(-?\d+)\s*\??', _unicode_replace, text)

    # RTF paragraph breaks
    text = re.sub(r'\\par\b', '\n', text)
    text = re.sub(r'\\par\s*', '\n', text)
    text = re.sub(r'\\line\s*', '\n', text)
    text = re.sub(r'\\tab\s*', '\t', text)

    # Remove remaining control words (\word, \word123, \*\word)
    text = re.sub(r'\\\*?\\[a-z]+\d*\s?', ' ', text, flags=re.IGNORECASE)

    # Remove braces and backslashes
    text = re.sub(r'[\\{}]', ' ', text)

    # Collapse whitespace and clean
    text = re.sub(r'\n\s*\n', '\n\n', text)
    text = re.sub(r' +', ' ', text)
    lines = [ln.strip() for ln in text.split('\n')]
    text = '\n'.join(ln for ln in lines if ln)

    if not text or len(text) < 10:
        return _extract_plain_text(data, '.rtf', 'RTF vacío o muy complejo')

    return ExtractionResult(
        pages=[{'page': 1, 'text': text, 'char_count': len(text)}],
        full_text=text,
        total_pages=1,
        total_chars=len(text),
        metadata=None,
    )


# ─── CSV / TSV ─────────────────────────────────────────────────────────────────

def _extract_csv_text(data: bytes, filename: str) -> ExtractionResult:
    ext = os.path.splitext(filename)[1].lower()
    delimiter = '\t' if ext == '.tsv' else ','
    text = data.decode('utf-8', errors='replace')
    rows = []
    reader = csv.reader(io.StringIO(text), delimiter=delimiter)
    for row in reader:
        rows.append(' | '.join(row).strip())
    full_text = '\n'.join(rows) if rows else text
    return ExtractionResult(
        pages=[{'page': 1, 'text': full_text, 'char_count': len(full_text)}],
        full_text=full_text,
        total_pages=1,
        total_chars=len(full_text),
        metadata=None,
    )


# ─── Plain text (code, markdown, HTML, JSON, YAML, etc.) ────────────────────────

def _extract_plain_text(data: bytes, filename: str = '', note: str = '') -> ExtractionResult:
    try:
        text = data.decode('utf-8', errors='replace')
    except Exception:
        text = data.decode('latin-1', errors='replace')
    text = text.strip()
    ext = os.path.splitext(filename)[1].lower()

    # HTML: strip tags to get readable content
    if ext in ('.html', '.htm'):
        text = _strip_html(text)

    if note:
        text = f'[{note}]\n\n{text}'

    return ExtractionResult(
        pages=[{'page': 1, 'text': text, 'char_count': len(text)}],
        full_text=text,
        total_pages=1,
        total_chars=len(text),
        metadata=None,
    )


def _strip_html(html_text: str) -> str:
    """Extract readable text from HTML using the stdlib HTMLParser."""
    from html.parser import HTMLParser

    class TextExtractor(HTMLParser):
        def __init__(self):
            super().__init__()
            self._parts = []
            self._skip = False
            self._skip_tag = ''

        def handle_starttag(self, tag, attrs):
            if tag in ('script', 'style'):
                self._skip = True
                self._skip_tag = tag

        def handle_endtag(self, tag):
            if tag == self._skip_tag and self._skip:
                self._skip = False
                self._skip_tag = ''
            # Block-level elements → newline
            if tag in ('p', 'div', 'li', 'br', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'tr', 'hr', 'section'):
                if self._parts and self._parts[-1] != '\n':
                    self._parts.append('\n')

        def handle_data(self, data):
            if not self._skip:
                t = data.strip()
                if t:
                    self._parts.append(t + ' ')

    extractor = TextExtractor()
    try:
        extractor.feed(html_text)
    except Exception:
        return html_text
    text = ''.join(extractor._parts)
    text = re.sub(r' {2,}', ' ', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


# ─── Archive (ZIP) ────────────────────────────────────────────────────────────

def _extract_archive(data: bytes, filename: str) -> ExtractionResult:
    """Extract text from all files inside a ZIP archive, concatenated."""
    parts = []
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            for name in zf.namelist():
                if name.endswith('/'):
                    continue
                try:
                    contents = zf.read(name)
                    inner_ext = os.path.splitext(name)[1].lower()
                    # Recurse into nested ZIPs (one level)
                    if inner_ext == '.zip':
                        inner = _extract_archive(contents, name)
                        parts.append(f'--- {name} ---\n{inner.full_text}')
                    elif inner_ext in PLAIN_TEXT_EXTENSIONS or inner_ext in STRUCTURED_TEXT_EXTENSIONS:
                        inner = _extract_plain_text(contents, name)
                        parts.append(f'--- {name} ---\n{inner.full_text}')
                    else:
                        # Try as plain text — may be garbled for binary
                        inner = _extract_plain_text(contents, name)
                        if inner.total_chars > 20:
                            parts.append(f'--- {name} ---\n{inner.full_text}')
                except Exception:
                    pass
    except Exception:
        return _extract_plain_text(data, filename, 'archivo no es un ZIP válido')

    if not parts:
        return _extract_plain_text(data, filename, 'sin contenido de texto extraíble')
    full_text = '\n\n'.join(parts)
    return ExtractionResult(
        pages=[{'page': 1, 'text': full_text, 'char_count': len(full_text)}],
        full_text=full_text,
        total_pages=1,
        total_chars=len(full_text),
        metadata=None,
    )


# ─── Backward-compat aliases ───────────────────────────────────────────────────

def extract_pdf(file_path: str) -> ExtractionResult:
    """Extract text from a PDF file on disk (legacy alias)."""
    with open(file_path, 'rb') as f:
        return _extract_pdf_bytes(f.read())


def extract_pdf_bytes(data: bytes) -> ExtractionResult:
    """Extract text from PDF bytes (legacy alias)."""
    return _extract_pdf_bytes(data)
